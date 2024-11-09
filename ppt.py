from pptx import Presentation
from pptx.util import Inches

# Create a PowerPoint presentation object
prs = Presentation()

# Add a blank slide layout
slide_layout = prs.slide_layouts[5]  # 5 corresponds to a blank slide layout
slide = prs.slides.add_slide(slide_layout)

# Add title
title_shape = slide.shapes.title
title_shape.text = "整机工艺流程与主体模块"

# Add process flow text box
left = Inches(0.5)
top = Inches(1.5)
width = Inches(4.5)
height = Inches(3)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "整机工艺流程：\n"
p.add_run().text = "1. 芯片/蓝膜剥离\n2. 芯片拾取/翻转\n3. 芯片绑定 + mark对准\n4. 位/力控键合\n5. 在线检测"

# Add corresponding modules text box
left = Inches(5)
top = Inches(1.5)
width = Inches(4.5)
height = Inches(3)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "对应主体模块：\n"
p.add_run().text = "1. 剥离模块\n2. 拾取翻转模块\n3. 视觉定位模块\n4. 纠姿键合模块\n5. 检测模块"

# Add auxiliary modules text box
left = Inches(0.5)
top = Inches(5)
width = Inches(9)
height = Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "辅助模块：\n"
p.add_run().text = "1. 洁净环境\n2. 防静电"

# Save the presentation
pptx_path = "整机工艺流程与主体模块.pptx"
prs.save(pptx_path)